from pptx import Presentation
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.util import Inches
from PIL import Image
import pprint
import os

def get_size(path_to_image):
    image = Image.open(path_to_image)
    dpi = image.info.get('dpi')

    return (image.width, image.height)

def set_header_of(a_slide, header_text):
    a_slide.placeholders[0].text_frame.text = header_text

def set_body_of(a_slide, body_text):
    a_slide.placeholders[1].text_frame.text = body_text

def add_picture_to(a_slide, image_path, slide_width):
    slide_body = a_slide.placeholders[1]
    resized_image_width, resized_image_height = resize_image(get_size(image_path), given_height=slide_body.height)

    top = slide_body.top
    left = slide_width / 2 - resized_image_width / 2

    a_slide.shapes.add_picture(image_path, left, top, height=resized_image_height)

def resize_image(original_size, given_height=None, given_width=None) -> (int, int):
    if given_height is not None and given_width is not None:
        raise "Only provide one argument"

    if given_height:
        resized_width = original_size[0] / original_size[1] * given_height
        return (resized_width, given_height)

    if given_width:
        resized_height = original_size[1] / original_size[0] * given_width
        return (given_width, resized_height)


def generate_presentation(incident_no: str, summary: str, session_id) -> Presentation:
    presentation: Presentation = Presentation()
    evidences_path = list(map(lambda evidence: f'./{session_id}/{evidence}', list(
    filter(lambda evidence: evidence.endswith(''), os.listdir(f'./{session_id}')))))

    first_slide = presentation.slides.add_slide(presentation.slide_layouts[0])

    set_header_of(first_slide, header_text="Decam report")
    set_body_of(first_slide, body_text=f"Incident Number {incident_no}")

    for evidence_path in evidences_path:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        set_header_of(slide, header_text=incident_no)
        add_picture_to(slide, image_path=evidence_path, slide_width=presentation.slide_width)

    stop_message_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    set_header_of(stop_message_slide, header_text="Stop Message")
    set_body_of(stop_message_slide, body_text=summary)

    
    return presentation